import os
import sys
import logging
from pathlib import Path
import cv2
import numpy as np
from PIL import Image, ImageDraw, ImageFont

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("deck_converter")

def render_text_paragraphs_to_slide_images(paragraphs: list, output_dir: str, doc_title: str = "DOCUMENT OVERVIEW") -> list:
    """Render a list of text paragraphs into slide images (e.g. from Word .docx or .txt)."""
    out_path = Path(output_dir)
    out_path.mkdir(parents=True, exist_ok=True)
    images = []

    width, height = 1280, 720
    bg_color = (15, 23, 42) # Dark navy
    accent_color = (37, 99, 235) # Cobalt blue

    # Group paragraphs into slides (approx 3-4 paragraphs per slide)
    chunk_size = 3
    chunks = [paragraphs[i:i + chunk_size] for i in range(0, len(paragraphs), chunk_size)]

    for page_num, chunk in enumerate(chunks, start=1):
        image = Image.new("RGB", (width, height), bg_color)
        draw = ImageDraw.Draw(image)

        try:
            font_title = ImageFont.truetype("arial.ttf", 38)
            font_sub = ImageFont.truetype("arial.ttf", 22)
            font_body = ImageFont.truetype("arial.ttf", 20)
        except Exception:
            font_title = font_sub = font_body = ImageFont.load_default()

        # Draw Header
        draw.rectangle([60, 50, width - 60, 56], fill=accent_color)
        draw.text((60, 70), f"{doc_title.upper()} · PAGE {page_num}", font=font_sub, fill=(148, 163, 184))
        draw.text((60, 110), f"Section {page_num}: Document Executive Summary", font=font_title, fill=(255, 255, 255))

        # Draw Body Card Container
        card_top = 200
        card_bottom = height - 60
        draw.rectangle([60, card_top, width - 60, card_bottom], fill=(30, 41, 59), outline=(51, 65, 85), width=2)

        y_pos = card_top + 30
        for p in chunk:
            p_text = p.strip()
            if not p_text:
                continue

            # Text wrapping
            max_chars = 70
            words = p_text.split()
            lines = []
            curr_line = ""
            for w in words:
                if len(curr_line) + len(w) + 1 > max_chars:
                    lines.append(curr_line)
                    curr_line = w
                else:
                    curr_line = (curr_line + " " + w).strip()
            if curr_line:
                lines.append(curr_line)

            for line in lines[:4]: # Limit lines per paragraph block
                draw.text((100, y_pos), line, font=font_body, fill=(241, 245, 249))
                y_pos += 28
            y_pos += 15 # Gap between paragraphs

            if y_pos > card_bottom - 40:
                break

        img_path = str(out_path / f"slide_{page_num:02d}.png")
        image.save(img_path)
        images.append(img_path)

    return images

def convert_docx_to_images(docx_path: str, output_dir: str) -> list:
    """Extract paragraphs from Word .docx file and convert to slide images."""
    try:
        import docx
        doc = docx.Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if not paragraphs:
            paragraphs = ["No readable text found in Word document."]
        title = Path(docx_path).stem.replace("_", " ").title()
        return render_text_paragraphs_to_slide_images(paragraphs, output_dir, doc_title=title)
    except Exception as e:
        logger.error(f"Error reading docx file {docx_path}: {e}")
        # Fallback raw text reader
        with open(docx_path, "rb") as f:
            raw_content = f.read().decode("latin1", errors="ignore")
        paragraphs = [line for line in raw_content.split("\n") if len(line.strip()) > 10][:12]
        return render_text_paragraphs_to_slide_images(paragraphs or ["Word Document Overview"], output_dir)

def convert_pdf_to_images(pdf_path: str, output_dir: str) -> list:
    """Convert PDF pages to high-res PNG image files using PyMuPDF/fitz or pypdf."""
    images = []
    out_path = Path(output_dir)
    out_path.mkdir(parents=True, exist_ok=True)
    
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(pdf_path)
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            img_name = str(out_path / f"slide_{i+1:02d}.png")
            pix.save(img_name)
            images.append(img_name)
    except Exception as e:
        logger.warning(f"fitz PDF extraction notice: {e}")

    if not images:
        try:
            import pypdf
            reader = pypdf.PdfReader(pdf_path)
            for i, page in enumerate(reader.pages):
                for count, image_file_object in enumerate(page.images):
                    img_name = str(out_path / f"slide_{i+1:02d}.png")
                    with open(img_name, "wb") as fp:
                        fp.write(image_file_object.data)
                    images.append(img_name)
                    break
        except Exception as e:
            logger.warning(f"PyPDF image extraction notice: {e}")

    return images

def extract_pdf_slides_and_text(pdf_path: str, output_dir: str) -> list:
    """Extract page text and render high-res slide images from PDF."""
    slide_items = []
    out_path = Path(output_dir)
    out_path.mkdir(parents=True, exist_ok=True)

    try:
        import fitz
        doc = fitz.open(pdf_path)
        for i, page in enumerate(doc):
            text = page.get_text().strip()
            pix = page.get_pixmap(dpi=150)
            img_name = str(out_path / f"slide_{i+1:02d}.png")
            pix.save(img_name)

            first_line = text.split("\n")[0] if text else f"Slide {i+1}"
            slide_items.append({
                "slide_number": i + 1,
                "title": first_line[:50].strip(),
                "text": text,
                "image_path": img_name
            })
        return slide_items
    except Exception as e:
        logger.warning(f"fitz PDF text extraction notice: {e}")

    # Fallback to convert_pdf_to_images
    img_paths = convert_pdf_to_images(pdf_path, output_dir)
    for i, path in enumerate(img_paths):
        slide_items.append({
            "slide_number": i + 1,
            "title": f"Slide {i+1}",
            "text": f"Slide {i+1} visual presentation",
            "image_path": path
        })
    return slide_items

def extract_pptx_slides_and_text(pptx_path: str, output_dir: str) -> list:
    """Extract native shape text, titles, speaker notes, and render slides from PPTX."""
    slide_items = []
    out_path = Path(output_dir)
    out_path.mkdir(parents=True, exist_ok=True)

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        doc_title = Path(pptx_path).stem.replace("_", " ").title()

        for i, slide in enumerate(prs.slides):
            text_runs = []
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    tf_text = shape.text_frame.text.strip()
                    if tf_text:
                        if not title and len(tf_text) < 80:
                            title = tf_text.split("\n")[0]
                        text_runs.append(tf_text)

            # Speaker notes if present
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    text_runs.append(f"Speaker Notes: {notes_text}")

            combined_text = "\n".join(text_runs) or f"Slide {i+1} Content"
            slide_title = title or (text_runs[0].split("\n")[0] if text_runs else f"Slide {i+1}")

            # Render synthetic slide card image from PPTX slide content
            img_path = render_single_slide_card(i + 1, slide_title, text_runs, out_path, doc_title)
            slide_items.append({
                "slide_number": i + 1,
                "title": slide_title[:50],
                "text": combined_text,
                "image_path": img_path
            })

        logger.info(f"Successfully extracted {len(slide_items)} slides from PPTX {Path(pptx_path).name}")
        return slide_items
    except Exception as e:
        logger.error(f"Error extracting PPTX slides: {e}")
        return []

def render_single_slide_card(slide_num: int, title: str, paragraphs: list, out_path: Path, doc_title: str) -> str:
    """Render crisp 1280x720 PNG slide image from extracted text blocks."""
    width, height = 1280, 720
    bg_color = (15, 23, 42)
    accent_color = (56, 189, 248) # Sky blue accent

    image = Image.new("RGB", (width, height), bg_color)
    draw = ImageDraw.Draw(image)

    try:
        font_title = ImageFont.truetype("arial.ttf", 36)
        font_sub = ImageFont.truetype("arial.ttf", 20)
        font_body = ImageFont.truetype("arial.ttf", 22)
    except Exception:
        font_title = font_sub = font_body = ImageFont.load_default()

    # Draw Header Accent Bar
    draw.rectangle([60, 45, width - 60, 51], fill=accent_color)
    draw.text((60, 65), f"{doc_title.upper()} · SLIDE {slide_num}", font=font_sub, fill=(148, 163, 184))
    draw.text((60, 100), title[:55], font=font_title, fill=(255, 255, 255))

    # Draw Body Container
    card_top = 180
    card_bottom = height - 50
    draw.rectangle([60, card_top, width - 60, card_bottom], fill=(30, 41, 59), outline=(51, 65, 85), width=2)

    y_pos = card_top + 30
    for p in paragraphs[:5]:
        p_clean = p.replace("\n", " ").strip()
        if not p_clean:
            continue
        
        words = p_clean.split()
        lines = []
        curr_line = ""
        for w in words:
            if len(curr_line) + len(w) + 1 > 65:
                lines.append(curr_line)
                curr_line = w
            else:
                curr_line = (curr_line + " " + w).strip()
        if curr_line:
            lines.append(curr_line)

        for line in lines[:3]:
            draw.text((90, y_pos), line, font=font_body, fill=(241, 245, 249))
            y_pos += 32
        y_pos += 12
        if y_pos > card_bottom - 40:
            break

    img_filename = str(out_path / f"slide_{slide_num:02d}.png")
    image.save(img_filename)
    return img_filename

def create_video_from_slide_images(image_paths: list, output_mp4_path: str, duration_per_slide: float = 6.0) -> str:
    """Convert a list of slide images into a video MP4 file."""
    if not image_paths:
        raise ValueError("No slide images available for video creation.")

    out_dir = Path(output_mp4_path).parent
    out_dir.mkdir(parents=True, exist_ok=True)

    first_img = cv2.imread(image_paths[0])
    if first_img is None:
        raise ValueError(f"Unable to read image: {image_paths[0]}")

    target_width, target_height = 1280, 720
    fps = 30
    fourcc = cv2.VideoWriter_fourcc(*'mp4v')
    out = cv2.VideoWriter(output_mp4_path, fourcc, fps, (target_width, target_height))

    for img_path in image_paths:
        frame = cv2.imread(img_path)
        if frame is None:
            continue
        resized_frame = cv2.resize(frame, (target_width, target_height))
        for _ in range(int(duration_per_slide * fps)):
            out.write(resized_frame)

    out.release()
    logger.info(f"Created slide presentation video MP4: {output_mp4_path}")
    return output_mp4_path
